from PyPDF2 import PdfReader
import docx
from pptx import Presentation

def extract_text_from_path(path: str) -> str:
    """
    Extracts and returns plain text from .pdf, .docx, or .pptx files.
    Raises ValueError for unsupported extensions.
    """
    path_lower = path.lower()
    text_chunks = []

    if path_lower.endswith(".pdf"):
        with open(path, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text_chunks.append(t)

    elif path_lower.endswith(".docx"):
        doc = docx.Document(path)
        for para in doc.paragraphs:
            if para.text:
                text_chunks.append(para.text)

    elif path_lower.endswith(".pptx"):
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_chunks.append(shape.text)
    else:
        raise ValueError("Unsupported file type. Use .pdf, .docx or .pptx")

    return "\n\n".join(text_chunks).strip()
